#!/usr/bin/env python3
"""
make_architecture_slides.py
────────────────────────────────────────────────────────────────
Generate a technical review deck of the agentic pipeline built on
2026-06-11: the MCP data layer, the agentic reception agent, the
planner, and the Tier-2 expander.

Content is written directly (no LLM) for technical accuracy. Reuses
the project's "Ocean Gradient" palette for visual consistency.

Usage:
    module load pytorch/2.8.0
    python3 tools/make_architecture_slides.py            # -> RCSFA_agentic_pipeline_review.pptx
    python3 tools/make_architecture_slides.py --output my.pptx
"""
import argparse

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette (Ocean Gradient) ──────────────────────────────────────
NAVY  = RGBColor(0x0B, 0x25, 0x45)
TEAL  = RGBColor(0x1C, 0x72, 0x93)
SKY   = RGBColor(0x21, 0xA6, 0xD3)
SLATE = RGBColor(0x3A, 0x5A, 0x7C)
BODY  = RGBColor(0x22, 0x33, 0x44)
GREEN = RGBColor(0x1B, 0x8A, 0x5A)
AMBER = RGBColor(0xB9, 0x7A, 0x12)
PALE  = RGBColor(0xF7, 0xF9, 0xFB)
PANEL = RGBColor(0xE8, 0xEE, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xD8, 0xE3, 0xED)

FONT_T = "Calibri"
FONT_B = "Calibri Light"
FONT_M = "Consolas"
SW, SH = Inches(13.33), Inches(7.5)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _text(slide, l, t, w, h, runs, size=14, color=BODY, bold=False,
          font=FONT_B, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, ov in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(ov.get("size", size)); r.font.name = ov.get("font", font)
        r.font.bold = ov.get("bold", bold); r.font.color.rgb = ov.get("color", color)
    return tb


def _header(slide, title, kicker=None):
    _rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    _rect(slide, 0, Inches(1.15), SW, Inches(0.06), SKY)
    if kicker:
        _text(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.3),
              kicker.upper(), size=12, color=SKY, bold=True, font=FONT_T)
    _text(slide, Inches(0.6), Inches(0.44), Inches(12.1), Inches(0.66),
          title, size=26, color=WHITE, bold=True, font=FONT_T,
          anchor=MSO_ANCHOR.MIDDLE)


def _split_lead(t):
    if t.startswith("**") and t.count("**") >= 2:
        end = t.index("**", 2)
        return t[2:end], t[end + 2:]
    return None, t


def _bullets(slide, items, left=Inches(0.7), top=Inches(1.5),
             width=Inches(12.0), size=14, gap=7):
    tb = slide.shapes.add_textbox(left, top, width, Inches(5.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        level, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level; p.space_after = Pt(gap)
        glyph = "•  " if level == 0 else "–  "
        sz = size - (1 if level else 0)
        lead, rest = _split_lead(text)
        if lead:
            r = p.add_run(); r.text = glyph + lead
            r.font.bold = True; r.font.color.rgb = TEAL
            r.font.size = Pt(sz); r.font.name = FONT_B
            r2 = p.add_run(); r2.text = rest
            r2.font.color.rgb = BODY; r2.font.size = Pt(sz); r2.font.name = FONT_B
        else:
            r = p.add_run(); r.text = glyph + text
            r.font.color.rgb = BODY; r.font.size = Pt(sz); r.font.name = FONT_B
    return tb


def _table(slide, headers, rows, left, top, width, col_w=None,
           fs=12, hfs=12):
    nr, nc = len(rows) + 1, len(headers)
    t = slide.shapes.add_table(nr, nc, left, top, width, Inches(0.4 * nr)).table
    if col_w:
        for j, w in enumerate(col_w):
            t.columns[j].width = w
    for j, h in enumerate(headers):
        c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = TEAL
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
        r.font.bold = True; r.font.size = Pt(hfs); r.font.color.rgb = WHITE
        r.font.name = FONT_T
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else PALE
            p = c.text_frame.paragraphs[0]; p.word_wrap = True
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(fs); r.font.color.rgb = BODY; r.font.name = FONT_B
    return t


def _mono(slide, lines, left, top, width, height, fs=11):
    box = _rect(slide, left, top, width, height, NAVY)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        r.font.name = FONT_M; r.font.size = Pt(fs)
        r.font.color.rgb = RGBColor(0xCF, 0xE8, 0xF5)
    return box


def _flowbox(slide, l, t, w, h, title, sub, fill, tcol=WHITE):
    _rect(slide, l, t, w, h, fill)
    _text(slide, l, t + Inches(0.08), w, Inches(0.34), title, size=14,
          color=tcol, bold=True, font=FONT_T, align=PP_ALIGN.CENTER)
    _text(slide, l, t + Inches(0.42), w, h - Inches(0.46), sub, size=11,
          color=tcol, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
def build(path):
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH

    # ── 1. TITLE ──────────────────────────────────────────────────
    s = _blank(prs)
    _rect(s, 0, 0, SW, SH, NAVY)
    _rect(s, 0, Inches(4.55), SW, Inches(0.06), SKY)
    _text(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4),
          "Agentic Multi-Agent Pipeline", size=44, color=WHITE, bold=True,
          font=FONT_T)
    _text(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.0),
          "MCP Data Layer · Reception Agent · Planner · Tier-2 Expander",
          size=22, color=SKY, font=FONT_B)
    _text(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.2),
          [("Technical review — RCSFA multi-agent framework\n", {"color": WHITE, "size": 15}),
           ("Built 2026-06-11  ·  natural-language request → concrete, data-grounded sampling columns  ·  dry (no execution)",
            {"color": RGBColor(0x9F, 0xC4, 0xDC), "size": 13})])

    # ── 2. ARCHITECTURE OVERVIEW ──────────────────────────────────
    s = _blank(prs); _header(s, "The three-tier pipeline", "Architecture")
    col = Inches(3.6); x = Inches(4.0); w = Inches(5.3); h = Inches(0.82)
    ys = [Inches(1.55), Inches(2.62), Inches(3.69), Inches(4.76), Inches(5.83)]
    _flowbox(s, x, ys[0], w, h, "User request (natural language)",
             "“…partitioning & drivers in the Naches sub-watershed, validate with obs”",
             PANEL, tcol=BODY)
    _flowbox(s, x, ys[1], w, h, "RECEPTION  (LLM agent + MCP tools)",
             "classify · pick archetype · gather data · reason → framed BRIEF", TEAL)
    _flowbox(s, x, ys[2], w, h, "PLANNER  (LLM)",
             "design strategy: sampling_plan · validation · requires_capabilities", TEAL)
    _flowbox(s, x, ys[3], w, h, "TIER-2 EXPANDER  (Python, deterministic)",
             "real DEM/soil/Fan → concrete (lat,lon) columns", SLATE)
    _flowbox(s, x, ys[4], w, h, "Tier-3 per-column ELM config  (existing builder)",
             "not yet wired to expander output", PANEL, tcol=SLATE)
    for y in ys[:-1]:
        _text(s, x, y + h - Inches(0.02), w, Inches(0.22), "▼", size=12,
              color=SKY, bold=True, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.5), Inches(1.7), Inches(3.2), Inches(4.0),
          [("Tier 1 — LLM\n", {"color": TEAL, "bold": True, "size": 14}),
           ("reasoning & design\n\n\n", {"color": SLATE, "size": 12}),
           ("Tier 2 — Python\n", {"color": SLATE, "bold": True, "size": 14}),
           ("deterministic geodata\n(no hallucinated coords)", {"color": SLATE, "size": 12})])
    _text(s, Inches(9.6), Inches(1.7), Inches(3.3), Inches(4.5),
          [("Design principle\n", {"color": NAVY, "bold": True, "size": 14}),
           ("LLM designs the STRATEGY; Python materializes it against real data; "
            "the strict per-column config (Tier 3) stays untouched — so open "
            "reasoning never corrupts a model setup.", {"color": BODY, "size": 12})])

    # ── 3. MCP LAYER ──────────────────────────────────────────────
    s = _blank(prs); _header(s, "MCP data layer — how it works", "Model Context Protocol")
    _bullets(s, [
        "**What it is:** each data source is a standalone MCP server (a small Python process) speaking the Model Context Protocol over stdio.",
        "**Config-driven:** mcp_config.json lists each server (command + args + timeout). Adding a source = one config block + one server file. No core changes.",
        "**Plumbing:** MCPManager loads the config → MCPClient per server. MCPClient opens a fresh stdio session per call (HPC/async-safe), returns parsed JSON.",
        "**Tool schema:** every server exposes typed tools via list_tools(); the agent layer turns these into OpenAI function schemas automatically.",
        "**Runtime:** servers + clients run under module load pytorch/2.8.0 (the env where openai + mcp live). All tools are READ-ONLY data fetches — low risk.",
        (1, "5 servers today: weather · geology · usgs_water · terrain · fan_wtd"),
    ], size=15, gap=9)

    # ── 4. MCP SERVERS TABLE ──────────────────────────────────────
    s = _blank(prs); _header(s, "The five MCP servers", "MCP data layer")
    _table(s,
           ["Server", "Source", "Key tools", "Shape"],
           [["weather", "NWS / Open-Meteo", "get_climate_summary", "point"],
            ["geology", "USDA SSURGO", "get_soil_profile, get_pflotran_materials", "point"],
            ["usgs_water", "USGS OGC API", "get_groundwater_sites, get_water_table_depth,\nget_monitoring_locations (streamflow)", "bbox"],
            ["terrain ★", "USGS 3DEP + WBD", "resolve_watershed, get_elevation,\nsample_elevation_grid, elevation_summary", "point+bbox"],
            ["fan_wtd ★", "Fan et al. 2013 (local NetCDF)", "get_fan_wtd, sample_fan_wtd, data_status", "point+bbox"]],
           Inches(0.5), Inches(1.55), Inches(12.3),
           col_w=[Inches(1.7), Inches(2.5), Inches(6.0), Inches(2.1)], fs=12)
    _text(s, Inches(0.5), Inches(6.5), Inches(12), Inches(0.6),
          [("★ built this session.  ", {"color": GREEN, "bold": True, "size": 12}),
           ("terrain + fan_wtd + the groundwater tools on usgs_water close the planner's "
            "top data gaps (DEM, watershed boundary, observed & equilibrium water table).",
            {"color": SLATE, "size": 12})])

    # ── 5. MCP TECHNICAL DECISIONS ────────────────────────────────
    s = _blank(prs); _header(s, "MCP — key technical decisions & gotchas", "MCP data layer")
    _bullets(s, [
        "**Observed WTD via the OGC API, not legacy NWIS:** waterservices.usgs.gov is unreachable from NERSC (SSL handshake timeout); api.waterdata.usgs.gov works. Depth-to-water = parameter 72019 in the field-measurements collection.",
        "**Fan WTD is a static-file server:** lazy-loads a NAMERICA NetCDF tile (auto-discovered). Handles a time dimension (annual squeeze / monthly mean), applies the land mask, and auto-detects sign — the tiles store WTD negative-below, so it returns a positive depth_to_water_m.",
        "**terrain endpoints:** 3DEP EPQS for point/grid elevation (threaded); WBD ArcGIS REST resolves a HUC code OR name → bbox + area (Naches → HUC 17030002, 2,860.6 km²).",
        "**Point vs bbox:** SSURGO/weather are point-only; watershed work needs bbox/grid sampling → terrain.sample_elevation_grid + fan.sample_fan_wtd are the workhorses.",
        "**Coverage is real-world messy:** a CONUS sweep (12 sites) surfaces gaps — e.g. wells with zero records, ocean/edge no-data — before you design a study there.",
    ], size=14, gap=10)

    # ── 6. AGENTIC SHIFT + TOOL LOOP ──────────────────────────────
    s = _blank(prs); _header(s, "Reception: from scripted bridge → agentic tool-user", "Reception agent")
    _bullets(s, [
        "**The shift:** earlier, Python mechanically pulled fixed MCP data and the LLM only extracted fields. Now reception is a pure-LLM agent that DRIVES the tools itself and reasons over the results.",
        "**Why:** a scripted bridge can't decide what's relevant (streamflow vs wells vs Fan), can't handle a purely conceptual question, and doesn't use the data to make decisions.",
        "**ToolLoopAgent (generic runtime):** exposes every MCP tool as an OpenAI function named server__tool; loops create(tools=…) → dispatch tool_calls via MCPClient → feed results back → repeat (capped rounds) → return final message + full trace.",
        "**No domain logic in Python** — the loop is pure relay; all 'what to fetch / when to stop' lives in the LLM.",
        "**Gateway rule (verified):** the PNNL endpoint supports function-calling on gpt-5.5, gemini-2.5-flash, claude-opus-4-8, claude-sonnet-4-5 — but tools= must be passed on EVERY request or Bedrock-routed Claude 400s.",
    ], size=14, gap=10)

    # ── 7. RECEPTION BEHAVIOR ─────────────────────────────────────
    s = _blank(prs); _header(s, "Reception: what the agent decides", "Reception agent")
    _bullets(s, [
        "**1 · Classify intent:** clarification_needed · analyze_existing · design.",
        "**2 · Pick archetype (for design):**",
        (1, "conceptual — mechanism / no real site → call few/no tools (maybe one representative climate)"),
        (1, "site — tied to a real place → use the tools to resolve the domain + inventory heterogeneity & observations"),
        "**3 · Gather proportionally:** resolve_watershed → elevation_summary → groundwater wells → streamflow gages → Fan prior → soil. Small limits; only what THIS question needs.",
        "**4 · Reason over results (the point):**",
        (1, "no wells → mark WTD validation unavailable, note streamflow instead"),
        (1, "flat terrain → drop elevation as a driver;  Fan very deep → water table decoupled from root zone"),
        "**Rule:** only state values a tool returned — never invent coordinates, counts, or records.",
        "**Output:** one framed BRIEF (design_archetype, domain, heterogeneity, observations_available/missing, gaps) — or a clarification.",
    ], size=13, gap=6)

    # ── 8. RECEPTION LIVE EXAMPLE ─────────────────────────────────
    s = _blank(prs); _header(s, "Reception — live behavior on two questions", "Reception agent")
    _rect(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(5.2), PALE, line=LINE)
    _text(s, Inches(0.7), Inches(1.62), Inches(5.7), Inches(0.4),
          "SITE  ·  Naches watershed", size=16, color=TEAL, bold=True, font=FONT_T)
    _bullets(s, [
        "6 rounds, **10 tool calls**",
        "resolve_watershed → elevation_summary → groundwater_sites(25) → streamflow gages → Fan sample(64) → soil → checked 3 wells",
        "**archetype = site**, domain = Naches 2,860.6 km²",
        "**obs available:** water table · Fan prior · streamflow",
        "**obs missing:** in-situ soil moisture · gridded soil",
    ], left=Inches(0.7), top=Inches(2.05), width=Inches(5.7), size=13, gap=8)
    _rect(s, Inches(6.85), Inches(1.5), Inches(6.0), Inches(5.2), PALE, line=LINE)
    _text(s, Inches(7.05), Inches(1.62), Inches(5.6), Inches(0.4),
          "CONCEPTUAL  ·  texture → recharge/ET", size=16, color=SLATE, bold=True, font=FONT_T)
    _bullets(s, [
        "2 rounds, **1 tool call** (one representative climate)",
        "**archetype = conceptual**, domain = null",
        "drivers = soil_texture · forcing",
        "→ routed to the planner to design a CONTROLLED experiment (no site, no spatial sampling)",
    ], left=Inches(7.05), top=Inches(2.05), width=Inches(5.6), size=13, gap=8)

    # ── 9. PLANNER ARCHETYPES ─────────────────────────────────────
    s = _blank(prs); _header(s, "Planner: one designer, two archetypes", "Planner agent")
    _table(s,
           ["", "Conceptual / controlled", "Site / domain"],
           [["Trigger", "mechanism, no real site", "real watershed + obs"],
            ["Design", "vary the factor under controlled\nconditions (synthetic soils,\none forcing), isolate mechanism",
             "stratified spatial sampling across\nreal heterogeneity + validation"],
            ["N logic", "= factor levels (small, 3–6)", "= strata × replication (data-driven)"],
            ["Data", "little / none", "elevation, soil, wells, streamflow, Fan"],
            ["Validation", "optional / vs expectation", "pinned to in-domain observations"]],
           Inches(0.5), Inches(1.55), Inches(12.3),
           col_w=[Inches(1.7), Inches(5.3), Inches(5.3)], fs=12)
    _text(s, Inches(0.5), Inches(6.35), Inches(12.2), Inches(0.8),
          [("Reception sets design_archetype; the planner branches on it. ", {"color": SLATE, "size": 13}),
           ("The existing production ELM planner (forcing × soil × substrate) is already a conceptual/controlled design — so the vocabulary existed.",
            {"color": SLATE, "size": 13})])

    # ── 10. PLANNER OUTPUT CONTRACT ───────────────────────────────
    s = _blank(prs); _header(s, "Planner — open reasoning, structured output", "Planner agent")
    _bullets(s, [
        "**Capability prompt:** open scientific reasoning at STRATEGY altitude — it does NOT enumerate per-column configs or invent coordinates (that is Tier 2/3).",
        "**N is justified, not blind:** the planner chooses the column count from the question + heterogeneity and explains it — a skillful small design beats a dense grid.",
        "**Validation pinned to data:** each target names the variable, source, in-domain availability, comparison, and metric — honest about gaps (e.g. no in-situ soil moisture).",
    ], size=14, gap=8, top=Inches(1.45))
    _mono(s, [
        "{  scientific_decomposition: { goals, key_processes, drivers_to_resolve },",
        "   model_choice:        { primary_model, design_archetype, coupling_needed },",
        "   sampling_strategy:   { approach, strata, n_exploratory, n_justification },",
        "   sampling_plan:       [ { group, n, reason } ],        # the n | group | why table",
        "   validation_design:   [ { target_variable, observation_source,",
        "                            in_domain_available, comparison, metric } ],",
        "   requires_capabilities: [ { capability, why, blocks } ]  }",
    ], Inches(0.6), Inches(4.0), Inches(12.1), Inches(2.7), fs=12)

    # ── 11. TIER-2 EXPANDER ───────────────────────────────────────
    s = _blank(prs); _header(s, "Tier-2 expander — strategy → concrete columns", "Deterministic")
    _bullets(s, [
        "**Deterministic Python, no LLM:** the planner's strategy + the real DEM/soil/Fan → concrete (lat,lon) columns. No hallucinated coordinates.",
        "**Algorithm:** sample DEM grid (terrain) → equal-interval elevation bands → allocate N ∝ occupied area (≥1/band) → farthest-point selection for spatial spread → enrich each column with Fan WTD + SSURGO texture.",
    ], size=13, gap=7, top=Inches(1.4), width=Inches(12.2))
    _text(s, Inches(0.6), Inches(3.05), Inches(12), Inches(0.35),
          "Naches result — 12 columns across 4 bands (306–1968 m):", size=13,
          color=TEAL, bold=True)
    _table(s,
           ["id", "lat", "lon", "elev_m", "band", "Fan WTD (m)", "soil"],
           [["col_01", "46.7249", "-120.7078", "447", "1", "0.70  (valley, shallow)", "loam"],
            ["col_06", "46.6268", "-121.2511", "996", "2", "0.00", "loam"],
            ["col_07", "46.8230", "-121.2511", "1511", "3", "9.78", "sandy loam"],
            ["col_09", "46.4796", "-120.9406", "1151", "3", "202.1  (ridge, deep)", "—"],
            ["col_11", "46.7249", "-121.3287", "1629", "4", "58.76", "sandy loam"]],
           Inches(0.6), Inches(3.5), Inches(12.1),
           col_w=[Inches(1.2), Inches(1.5), Inches(1.7), Inches(1.3),
                  Inches(0.9), Inches(3.6), Inches(1.9)], fs=12)
    _text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.5),
          "Valley water table 0.7 m vs ridge 202 m — physically coherent; columns.json saved per run.",
          size=12, color=SLATE)

    # ── 12. ENGINEERING & TESTING ─────────────────────────────────
    s = _blank(prs); _header(s, "Engineering, safety & testing", "Quality")
    _bullets(s, [
        "**Additive & safe:** the legacy two-pass reception and production workflow.py are untouched; new agents are a parallel path.",
        "**Dry by design:** the whole pipeline stops at planned columns — no experiment manager, no SLURM. Transcripts saved per run (brief, tool trace, plan, columns).",
        "**Tests — 62 passing** under pytorch/2.8.0:",
        (1, "36 guardrail (existing ELM) · 20 MCP + expander (band/allocate/parse/fan-sign) · 6 agentic (tool-schema gen + brief parse)"),
        "**CONUS sweep:** tools/mcp_conus_sweep.py — coverage diagnostic with an --assert sanity gate.",
        "**New files:** mcp/terrain-mcp, mcp/fan-wtd-mcp, usgs groundwater_api; src/agents/tool_loop.py, reception_llm.py, prompts; tools/run_pipeline.py, expand_sampling.py.",
    ], size=14, gap=9)

    # ── 13. ROADMAP ───────────────────────────────────────────────
    s = _blank(prs); _header(s, "What's next", "Roadmap")
    _bullets(s, [
        "**Agentic analyzer (next):** give the analyzer observation-retrieval tools so it FETCHES the wells/streamflow/Fan and compares them to ELM output — where 'validate against observations' actually happens.",
        "**Wire Tier-2 → Tier-3:** feed columns.json into the existing per-column ELM config + run path.",
        "**Gridded soil (GSDE):** add the BNU global soil dataset as a fan_wtd-style MCP to make soil a real stratification axis (deferred today).",
        "**Sampling refinement:** topographic position (valley vs hillslope) via TWI / height-above-drainage — beyond elevation bands.",
        "**Still open:** FLUXNET soil-moisture ingest · ELM→PFLOTRAN weak coupling.",
    ], size=15, gap=11)
    _text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.5),
          "Principle going forward: tool-use is a shared capability — grant each agent the subset its job needs (including none).",
          size=12, color=SLATE)

    prs.save(path)
    return len(prs.slides._sldIdLst)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--output", default="RCSFA_agentic_pipeline_review.pptx")
    args = ap.parse_args()
    n = build(args.output)
    print(f"Saved {n}-slide deck -> {args.output}")


if __name__ == "__main__":
    main()
