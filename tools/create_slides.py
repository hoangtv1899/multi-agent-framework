#!/usr/bin/env python3
"""
create_slides.py
────────────────────────────────────────────────────────────────
Generate a polished summary deck from PFLOTRAN workflow results.
Uses Claude (via OpenAI-compatible API) to write slide content,
then assembles a 5-slide .pptx with the "Ocean Gradient" template.

Usage:
    python3 tools/create_slides.py \
        --run-dir workflow_outputs/run_20260331_194240 \
        --output  summary_slides.pptx

Requirements:
    pip install --user python-pptx Pillow openai
    Environment variable: PNNL_API_KEY
────────────────────────────────────────────────────────────────
"""

import os
import re
import json
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ═══════════════════════════════════════════════════════════════
#  DESIGN TOKENS — Ocean Gradient palette
# ═══════════════════════════════════════════════════════════════

NAVY       = RGBColor(0x0B, 0x25, 0x45)
TEAL       = RGBColor(0x1C, 0x72, 0x93)
SKY        = RGBColor(0x21, 0xA6, 0xD3)
LIGHT_BLUE = RGBColor(0x65, 0xD2, 0xE9)
MUTED_BLUE = RGBColor(0x8A, 0xBE, 0xDB)
SLATE      = RGBColor(0x3A, 0x5A, 0x7C)
PALE_BG    = RGBColor(0xF7, 0xF9, 0xFB)
PANEL_BG   = RGBColor(0xE8, 0xEE, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER= RGBColor(0xD8, 0xE3, 0xED)
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2E)
DIVIDER    = RGBColor(0xE2, 0xEA, 0xF0)

FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri Light"
SLIDE_W      = Inches(13.33)
SLIDE_H      = Inches(7.5)

# Minimum font sizes — all body text >= 12pt per user request
BODY_SIZE    = Pt(14)
BULLET_SIZE  = Pt(14)
CAPTION_SIZE = Pt(12)
LABEL_SIZE   = Pt(10)
TITLE_SIZE   = Pt(26)


# ═══════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════

def _set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text, left, top, width, height,
                 font_name=FONT_BODY, font_size=BODY_SIZE,
                 color=NAVY, bold=False, alignment=PP_ALIGN.LEFT,
                 word_wrap=True, line_spacing=None):
    """Add a single-run text box and return the shape."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def _add_rich_text(slide, segments, left, top, width, height,
                   alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add a text box with multiple styled runs on one line.
    segments: list of (text, font_name, font_size, color, bold)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    for text, fn, fs, col, b in segments:
        run = p.add_run()
        run.text = text
        run.font.name = fn
        run.font.size = fs
        run.font.color.rgb = col
        run.font.bold = b
    return txBox


def _add_bullets(slide, items, left, top, width, height,
                 font_name=FONT_BODY, font_size=BULLET_SIZE,
                 color=SLATE, bullet_color=TEAL,
                 line_spacing=Pt(24)):
    """Add a text box with bullet items (using · prefix)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.line_spacing = line_spacing

        # Bullet marker
        run_bullet = p.add_run()
        run_bullet.text = "·  "
        run_bullet.font.name = font_name
        run_bullet.font.size = Pt(font_size.pt + 4)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.bold = True

        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.name = font_name
        run_text.font.size = font_size
        run_text.font.color.rgb = color
    return txBox


def _add_section_label(slide, text, left, top, width,
                       color=TEAL, font_size=LABEL_SIZE):
    """Add an uppercase section label (e.g. LOCATION, CASES)."""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text.upper()
    run.font.name = FONT_TITLE
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = True
    return txBox


def _add_rect(slide, left, top, width, height, fill_color,
              border_color=None, border_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _add_rounded_rect(slide, left, top, width, height, fill_color,
                      border_color=None, border_width=None):
    """Add a filled rounded rectangle shape."""
    shape = _add_rect(slide, left, top, width, height, fill_color,
                      border_color, border_width)
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
    return shape


def _add_stat(slide, value, label, left, top,
              val_color=TEAL, lbl_color=MUTED_BLUE):
    """Add a large stat number with a small label below."""
    _add_textbox(slide, value,
                 left, top, Inches(2.2), Inches(0.6),
                 font_name=FONT_TITLE, font_size=Pt(36),
                 color=val_color, bold=True)
    _add_textbox(slide, label,
                 left, top + Inches(0.55), Inches(2.2), Inches(0.35),
                 font_name=FONT_BODY, font_size=CAPTION_SIZE,
                 color=lbl_color)


def _add_kv_row(slide, key, value, left, top, width,
                key_color=MUTED_BLUE, val_color=NAVY):
    """Add a key-value row with a subtle bottom border."""
    _add_textbox(slide, key,
                 left, top, Inches(3.5), Inches(0.35),
                 font_size=BODY_SIZE, color=key_color)
    _add_textbox(slide, value,
                 left + Inches(3.5), top, width - Inches(3.5), Inches(0.35),
                 font_size=BODY_SIZE, color=val_color, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top + Inches(0.38), width, Pt(0.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = DIVIDER
    line.line.fill.background()


def _try_add_image(slide, fig_path, left, top, width, height):
    """Add an image if the file exists; return True/False."""
    if fig_path and fig_path.exists():
        slide.shapes.add_picture(
            str(fig_path), left, top, width, height
        )
        return True
    return False


def _add_placeholder_text(slide, text, left, top, width, height):
    """Show a placeholder label where a figure would go."""
    _add_textbox(slide, f"[{text}]",
                 left, top, width, height,
                 font_size=CAPTION_SIZE, color=MUTED_BLUE,
                 alignment=PP_ALIGN.CENTER)


def _add_gradient_bar(slide, left, top, width, height=Pt(5)):
    """Add a thin horizontal gradient-like accent bar (3 segments)."""
    seg_w = width // 3
    for i, c in enumerate([TEAL, SKY, LIGHT_BLUE]):
        _add_rect(slide, left + seg_w * i, top, seg_w, height, c)


def _add_slide_number(slide, number, color=MUTED_BLUE):
    """Add slide number at bottom-right."""
    _add_textbox(slide, f"{number:02d}",
                 SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.5),
                 Inches(0.5), Inches(0.3),
                 font_size=CAPTION_SIZE, color=color,
                 alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_slide_1_title(prs, data):
    """Slide 1 — Dark title card."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, NAVY)

    # Top accent bar
    _add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Pt(5))

    # Title
    title = data.get("title", "PFLOTRAN Simulation Results")
    _add_textbox(slide, title,
                 Inches(1.0), Inches(2.0), Inches(10.0), Inches(1.6),
                 font_name=FONT_TITLE, font_size=Pt(38),
                 color=WHITE, bold=True)

    # Subtitle
    subtitle = data.get("subtitle", "")
    _add_textbox(slide, subtitle,
                 Inches(1.0), Inches(3.8), Inches(9.0), Inches(1.0),
                 font_name=FONT_BODY, font_size=Pt(18),
                 color=MUTED_BLUE)

    # Bottom-left: branding
    _add_rich_text(slide, [
        ("●  ", FONT_TITLE, Pt(12), TEAL, False),
        ("PFLOTRAN WORKFLOW", FONT_TITLE, LABEL_SIZE, MUTED_BLUE, False),
    ], Inches(1.0), SLIDE_H - Inches(0.7), Inches(3.0), Inches(0.3))

    # Bottom-right: tag
    tag = data.get("tag", "")
    if tag:
        _add_textbox(slide, tag,
                     SLIDE_W - Inches(3.5), SLIDE_H - Inches(0.7),
                     Inches(3.0), Inches(0.3),
                     font_size=CAPTION_SIZE, color=MUTED_BLUE,
                     alignment=PP_ALIGN.RIGHT)

    notes = data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_slide_2_setup(prs, data, figures):
    """Slide 2 — Experiment setup: text left, TWO figures stacked right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PALE_BG)

    # Right panel background
    panel_x = Inches(6.0)
    _add_rect(slide, panel_x, Inches(0), SLIDE_W - panel_x, SLIDE_H, PANEL_BG)

    # Title
    title = data.get("title", "Experiment Setup")
    _add_textbox(slide, title,
                 Inches(0.8), Inches(0.4), Inches(5.0), Inches(0.9),
                 font_name=FONT_TITLE, font_size=TITLE_SIZE,
                 color=NAVY, bold=True)

    # Subheaders + bullets — use Pt(13) for dense left panel
    y = Inches(1.5)
    subheaders = data.get("subheaders", [])
    left_font = Pt(13)
    for sh in subheaders:
        label = sh.get("label", "")
        bullets = sh.get("bullets", [])

        _add_section_label(slide, label, Inches(0.8), y, Inches(5.0))
        y += Inches(0.28)

        if bullets:
            n = len(bullets)
            box_h = Inches(0.30 * n + 0.1)
            _add_bullets(slide, bullets,
                         Inches(0.8), y, Inches(4.9), box_h,
                         color=SLATE, bullet_color=TEAL,
                         font_size=left_font, line_spacing=Pt(20))
            y += box_h + Inches(0.08)

    # TWO figures stacked on the right panel
    fig_left = Inches(6.3)
    fig_w = Inches(6.5)
    fig_h = Inches(3.2)

    # Top figure: comparison_experiments
    fig1 = figures.get("comparison_experiments")
    _add_textbox(slide, "Experiment comparison",
                 fig_left, Inches(0.3), fig_w, Inches(0.3),
                 font_size=CAPTION_SIZE, color=SLATE, bold=True)
    if not _try_add_image(slide, fig1,
                          fig_left, Inches(0.6), fig_w, fig_h):
        _add_placeholder_text(slide, "comparison_experiments.png",
                              fig_left, Inches(1.5), fig_w, Inches(0.5))

    # Bottom figure: comparison_flux
    fig2 = figures.get("comparison_flux")
    _add_textbox(slide, "Flux boundary conditions",
                 fig_left, Inches(3.95), fig_w, Inches(0.3),
                 font_size=CAPTION_SIZE, color=SLATE, bold=True)
    if not _try_add_image(slide, fig2,
                          fig_left, Inches(4.25), fig_w, fig_h):
        _add_placeholder_text(slide, "comparison_flux_conditions.png",
                              fig_left, Inches(5.2), fig_w, Inches(0.5))

    _add_slide_number(slide, 2)

    notes = data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_slide_3_results(prs, data, figures):
    """Slide 3 — Key results: stats + section bullets left, figure right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    # Title (should include recharge condition from LLM)
    title = data.get("title", "Key Results")
    _add_textbox(slide, title,
                 Inches(0.8), Inches(0.3), Inches(6.0), Inches(0.9),
                 font_name=FONT_TITLE, font_size=TITLE_SIZE,
                 color=NAVY, bold=True)

    # Parse subheaders — extract leading "number — label" as stats
    subheaders = data.get("subheaders", [])
    stats = []
    remaining_sections = []

    for sh in subheaders:
        label = sh.get("label", "")
        bullets = sh.get("bullets", [])
        section_bullets = []
        for b in bullets:
            m = re.match(r'^([0-9.]+\s*\S+)\s*[—–-]\s*(.+)$', b)
            if m and len(stats) < 3:
                stats.append((m.group(1).strip(), m.group(2).strip()))
            else:
                section_bullets.append(b)
        if section_bullets:
            remaining_sections.append((label, section_bullets))

    # If no stats parsed, just use all bullets directly
    if not stats:
        remaining_sections = []
        for sh in subheaders:
            remaining_sections.append(
                (sh.get("label", ""), sh.get("bullets", []))
            )

    # Stat callouts (top row)
    y = Inches(1.3)
    if stats:
        x = Inches(0.8)
        for val, lbl in stats[:3]:
            _add_stat(slide, val, lbl, x, y)
            x += Inches(2.2)
        y = Inches(2.5)

    # Section labels + bullets for remaining content
    for label, bullets in remaining_sections:
        if label:
            _add_section_label(slide, label, Inches(0.8), y, Inches(6.0))
            y += Inches(0.28)
        if bullets:
            n = len(bullets)
            box_h = Inches(0.28 * n + 0.1)
            _add_bullets(slide, bullets,
                         Inches(0.8), y, Inches(5.8), box_h,
                         color=SLATE, bullet_color=TEAL,
                         font_size=Pt(13), line_spacing=Pt(19))
            y += box_h + Inches(0.08)

    # Summary text
    summary = data.get("summary", "")
    if summary and y < Inches(6.5):
        _add_textbox(slide, summary,
                     Inches(0.8), y + Inches(0.1),
                     Inches(5.8), Inches(7.2) - y,
                     font_size=Pt(12), color=SLATE)

    # Figure (right side)
    fig = figures.get("saturation_comparison")
    if not _try_add_image(slide, fig,
                          Inches(7.0), Inches(0.4),
                          Inches(5.8), Inches(6.6)):
        _add_placeholder_text(slide, "comparison_all_times.png",
                              Inches(8.0), Inches(3.5), Inches(4.0), Inches(0.5))

    _add_slide_number(slide, 3)

    notes = data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_slide_4_performance(prs, data, figures):
    """Slide 4 — Solver performance: dark background, text + figure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BG)

    # Title
    title = data.get("title", "Solver Performance")
    _add_textbox(slide, title,
                 Inches(0.8), Inches(0.4), Inches(5.5), Inches(0.9),
                 font_name=FONT_TITLE, font_size=TITLE_SIZE,
                 color=WHITE, bold=True)

    # Subheaders + bullets — tighter spacing for dark slide
    y = Inches(1.4)
    subheaders = data.get("subheaders", [])
    for sh in subheaders:
        label = sh.get("label", "")
        bullets = sh.get("bullets", [])

        _add_section_label(slide, label, Inches(0.8), y, Inches(5.5),
                           color=SKY)
        y += Inches(0.30)

        if bullets:
            n = len(bullets)
            box_h = Inches(0.28 * n + 0.1)
            _add_bullets(slide, bullets,
                         Inches(0.8), y, Inches(5.5), box_h,
                         color=MUTED_BLUE, bullet_color=SKY,
                         font_size=Pt(13), line_spacing=Pt(20))
            y += box_h + Inches(0.06)

    # Summary text — only if space remains
    summary = data.get("summary", "")
    if summary and y < Inches(6.2):
        remaining = Inches(7.0) - y
        _add_textbox(slide, summary,
                     Inches(0.8), y + Inches(0.15), Inches(5.5), remaining,
                     font_size=Pt(13), color=MUTED_BLUE)

    # Figure (right half)
    fig = figures.get("performance")
    if not _try_add_image(slide, fig,
                          Inches(6.8), Inches(0.4),
                          Inches(6.0), Inches(6.6)):
        _add_placeholder_text(slide, "performance_summary.png",
                              Inches(8.0), Inches(3.5), Inches(4.0), Inches(0.5))

    _add_slide_number(slide, 4, color=MUTED_BLUE)

    notes = data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_slide_5_next_steps(prs, data):
    """Slide 5 — Next steps: 3 expanded numbered cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, PALE_BG)

    # Title
    title = data.get("title", "Next Steps")
    _add_textbox(slide, title,
                 Inches(0.8), Inches(0.4), Inches(10.0), Inches(0.8),
                 font_name=FONT_TITLE, font_size=TITLE_SIZE,
                 color=NAVY, bold=True)

    # Build cards from subheaders — wider and taller
    subheaders = data.get("subheaders", [])
    n_cards = min(len(subheaders), 3) or 3

    total_w = Inches(11.73)  # 13.33 - 2*0.8 margin
    gap = Inches(0.4)
    card_w = (total_w - gap * (n_cards - 1)) / n_cards
    card_h = Inches(4.8)
    start_x = Inches(0.8)
    card_y = Inches(1.5)

    for i, sh in enumerate(subheaders[:3]):
        x = start_x + i * (card_w + gap)
        label = sh.get("label", f"Step {i+1}")
        bullets = sh.get("bullets", [])

        # Card background
        _add_rounded_rect(slide, x, card_y, card_w, card_h, WHITE,
                          border_color=CARD_BORDER, border_width=Pt(0.75))

        # Large number
        _add_textbox(slide, f"{i+1:02d}",
                     x + Inches(0.35), card_y + Inches(0.25),
                     Inches(1.0), Inches(0.6),
                     font_name=FONT_TITLE, font_size=Pt(36),
                     color=TEAL, bold=True)

        # Card title
        _add_textbox(slide, label,
                     x + Inches(0.35), card_y + Inches(0.95),
                     card_w - Inches(0.7), Inches(0.45),
                     font_name=FONT_TITLE, font_size=Pt(16),
                     color=NAVY, bold=True)

        # Card body — larger font, fills card
        body = "\n".join(bullets) if bullets else ""
        _add_textbox(slide, body,
                     x + Inches(0.35), card_y + Inches(1.5),
                     card_w - Inches(0.7), card_h - Inches(1.9),
                     font_name=FONT_BODY, font_size=BODY_SIZE,
                     color=SLATE)

    # Footer bar
    footer_y = SLIDE_H - Inches(0.7)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), footer_y, total_w, Pt(0.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = CARD_BORDER
    line.line.fill.background()

    _add_textbox(slide, "PFLOTRAN automated workflow",
                 Inches(0.8), footer_y + Inches(0.1),
                 Inches(4.0), Inches(0.3),
                 font_size=CAPTION_SIZE, color=MUTED_BLUE)

    _add_slide_number(slide, 5)

    notes = data.get("speaker_notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ═══════════════════════════════════════════════════════════════
#  DATA LOADING
# ═══════════════════════════════════════════════════════════════

def load_run_data(run_dir: Path) -> dict:
    """Load analysis report and locate figures from run directory."""
    analysis_path = run_dir / "ANALYSIS_REPORT.json"
    if not analysis_path.exists():
        raise FileNotFoundError(f"Analysis report not found: {analysis_path}")

    with open(analysis_path) as f:
        analysis = json.load(f)

    figures = {
        "comparison_experiments": run_dir / "02_setup_plots" /
                                  "comparison_experiments.png",
        "comparison_flux":        run_dir / "02_setup_plots" /
                                  "comparison_flux_conditions.png",
        "saturation_comparison":  run_dir / "04_analysis" /
                                  "comparison_all_times.png",
        "performance":            run_dir / "03_results" /
                                  "performance_summary.png",
    }

    print(f"\n📂 Run directory: {run_dir.name}")
    print("\n📊 Figures found:")
    for name, p in figures.items():
        status = "✅" if p.exists() else "❌ missing"
        print(f"   {status}  {name}: {p.name}")

    return {"analysis": analysis, "figures": figures}


# ═══════════════════════════════════════════════════════════════
#  LLM CONTENT GENERATION
# ═══════════════════════════════════════════════════════════════

def write_slide_content(analysis: dict) -> dict:
    """Ask Claude to write concise slide content. Returns structured dict."""
    import openai

    print("\n🤖 Asking Claude to write slide content...")

    prompt = (
        f"Convert this PFLOTRAN analysis into 5 clean professional slides.\n\n"
        f"ANALYSIS:\n{json.dumps(analysis, indent=2)}\n\n"

        f"STRICT RULES:\n"
        f"  - Max 10 words per bullet\n"
        f"  - Lead every bullet with a number or symbol when possible\n"
        f"  - Group bullets under sub-headers\n"
        f"  - Slide titles must state the conclusion, not the topic\n"
        f"    BAD: 'Key Results'  GOOD: 'All Cases Converge — But Path Differs'\n"
        f"  - No full sentences — fragments only\n"
        f"  - No jargon — plain scientific English\n"
        f"  - Do NOT say 'surprising' or 'unexpected'\n"
        f"  - Be specific with numbers: include units, depths, times\n"
        f"  - ALWAYS include recharge rate (mm/yr) when naming the experiment\n"
        f"  - Simulation time in results must match actual data\n\n"

        f"SLIDE STRUCTURE — follow exactly:\n\n"

        f"  Slide 1 — TITLE\n"
        f"    title: experiment name + location + recharge rate (max 12 words)\n"
        f"    subtitle: summary of main finding + experiment count + sim length\n"
        f"    tag: date or run identifier\n\n"

        f"  Slide 2 — EXPERIMENT SETUP\n"
        f"    title: what was tested + why (include recharge rate)\n"
        f"    subheaders (3+ bullets per section to fill space):\n"
        f"      LOCATION: site name, climate zone, annual precipitation,\n"
        f"                recharge rate with units\n"
        f"      CASES: number of experiments, parameter varied,\n"
        f"             range of values tested, rationale\n"
        f"      GRID: cell count, total depth, vertical resolution,\n"
        f"            soil database source\n"
        f"      SIMULATION: total time, timestep strategy, output intervals\n\n"

        f"  Slide 3 — KEY RESULTS\n"
        f"    title: conclusion-first, MUST include recharge condition\n"
        f"           e.g. 'All Cases Converge by 2 Years at 63 mm/yr Recharge'\n"
        f"    subheaders (2-3 bullets per section for detail):\n"
        f"      SHALLOW ZONE (<1.0 m): equilibrium saturation value,\n"
        f"        time to equilibrium, sensitivity to initial WT depth\n"
        f"      DEEP ZONE (>2.0 m): initial saturation range,\n"
        f"        recovery trajectory, final equilibrium\n"
        f"      CONVERGENCE: when all cases merge, recharge rate influence,\n"
        f"        steady-state saturation profile\n"
        f"      CRITICAL TRANSITION: depth range, physical mechanism,\n"
        f"        capillary fringe behavior, why this zone matters\n"
        f"    summary: 1-2 sentence plain-English interpretation of results\n\n"

        f"  Slide 4 — SOLVER PERFORMANCE\n"
        f"    title: conclusion-first (e.g. 'Robust Solver — Zero Failures')\n"
        f"    subheaders (2-3 bullets per section):\n"
        f"      NEWTON ITERATIONS: range from easiest to hardest case,\n"
        f"        percent increase, which case needed most iterations\n"
        f"      WALL CLOCK TIME: time per case, total wall time,\n"
        f"        comparison to simulation length\n"
        f"      TIMESTEP BEHAVIOR: any cuts or reductions, adaptive strategy,\n"
        f"        convergence quality\n"
        f"      STABILITY: boundary condition quality, mass balance,\n"
        f"        overall solver reliability assessment\n"
        f"    summary: 1 sentence overall performance assessment\n\n"

        f"  Slide 5 — NEXT STEPS\n"
        f"    title: action-oriented (e.g. 'Three Paths Forward')\n"
        f"    subheaders (exactly 3, with 2-3 bullets each for detail):\n"
        f"      REFINE TRANSITION ZONE: specific depth range to target,\n"
        f"        number of new cases, expected insight, grid refinement\n"
        f"      OPTIMIZE RUNTIME: current sim length vs needed,\n"
        f"        potential speedup, resource savings\n"
        f"      NEXT VARIABLE: what to vary (e.g. recharge rate, soil type),\n"
        f"        scientific rationale, expected range, number of new runs\n\n"

        f"Output JSON only — no markdown fences, no comments:\n"
        f'{{"slides": [\n'
        f'  {{"id": 1, "title": "...", "subtitle": "...", "tag": "...", '
        f'"speaker_notes": "..."}},\n'
        f'  {{"id": 2, "title": "...", '
        f'"subheaders": [{{"label": "...", "bullets": ["...", "..."]}}], '
        f'"speaker_notes": "..."}},\n'
        f'  {{"id": 3, "title": "...", '
        f'"subheaders": [{{"label": "...", "bullets": ["...", "..."]}}], '
        f'"summary": "...", "speaker_notes": "..."}},\n'
        f'  {{"id": 4, "title": "...", '
        f'"subheaders": [{{"label": "...", "bullets": ["...", "..."]}}], '
        f'"summary": "...", "speaker_notes": "..."}},\n'
        f'  {{"id": 5, "title": "...", '
        f'"subheaders": [{{"label": "...", "bullets": ["...", "..."]}}], '
        f'"speaker_notes": "..."}}\n'
        f']}}\n\n'
        f"CRITICAL JSON RULES:\n"
        f"  - Output ONLY valid JSON — no markdown fences\n"
        f"  - No trailing commas after last item in array or object\n"
        f"  - Use spaces for newlines inside strings — never raw newlines\n"
        f"  - Escape any quotes inside strings with \\\"\n"
        f"  - Max 10 words per bullet to reduce JSON errors\n"
    )

    client = openai.OpenAI(
        api_key  = os.getenv("PNNL_API_KEY"),
        base_url = "https://ai-incubator-api.pnnl.gov",
    )

    response = client.chat.completions.create(
        model    = "claude-opus-4-5-20251101-v1-project",
        messages = [
            {"role": "system",
             "content": (
                 "You are a scientific presentation writer. "
                 "Output ONLY valid JSON — no markdown fences, "
                 "no comments, no trailing commas. "
                 "Be specific with numbers and units. "
                 "Always include recharge rate when naming experiments."
             )},
            {"role": "user", "content": prompt}
        ]
    )

    raw = response.choices[0].message.content
    return _parse_llm_json(raw, client)


def _parse_llm_json(raw: str, client=None) -> dict:
    """Robustly parse JSON from LLM output, with retry."""
    raw = re.sub(r"```json\s*", "", raw)
    raw = re.sub(r"```\s*",     "", raw)
    raw = raw.strip()

    start = raw.find("{")
    end   = raw.rfind("}")
    if start == -1 or end == -1:
        raise ValueError("No JSON object found in response")

    json_str = raw[start:end + 1]

    json_str = re.sub(r",\s*}", "}", json_str)
    json_str = re.sub(r",\s*]", "]", json_str)
    json_str = re.sub(r'(?<!\\)\n', ' ', json_str)

    try:
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        print(f"   ⚠️  JSON parse failed: {e}")
        print(f"   Near: {json_str[max(0, e.pos-50):e.pos+50]}")

        if client:
            print("   🔄 Asking Claude to fix JSON...")
            fix_response = client.chat.completions.create(
                model    = "claude-opus-4-5-20251101-v1-project",
                messages = [
                    {"role": "system",
                     "content": "Fix the JSON syntax error. "
                                "Output ONLY valid JSON."},
                    {"role": "user",
                     "content": f"Fix this invalid JSON:\n{json_str}"}
                ]
            )
            fixed = fix_response.choices[0].message.content.strip()
            fixed = re.sub(r"```json\s*", "", fixed)
            fixed = re.sub(r"```\s*",     "", fixed).strip()
            s = fixed.find("{")
            e2 = fixed.rfind("}")
            return json.loads(fixed[s:e2 + 1])
        raise


# ═══════════════════════════════════════════════════════════════
#  ASSEMBLY
# ═══════════════════════════════════════════════════════════════

def assemble_slides(slide_content: dict, figures: dict, output_path: str):
    """Build the .pptx from structured content + figures."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides_data = slide_content.get("slides", [])

    builders = {
        1: build_slide_1_title,
        2: build_slide_2_setup,
        3: build_slide_3_results,
        4: build_slide_4_performance,
        5: build_slide_5_next_steps,
    }

    for sd in slides_data:
        sid = sd.get("id", 0)
        builder = builders.get(sid)
        if not builder:
            print(f"   ⚠️  Unknown slide id {sid}, skipping")
            continue

        print(f"   🎨 Building slide {sid}: {sd.get('title', '?')[:50]}")

        if sid in (2, 3, 4):
            builder(prs, sd, figures)
        else:
            builder(prs, sd)

    prs.save(output_path)
    print(f"\n✅ Slides saved to: {output_path}")


# ═══════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="Generate summary slides from PFLOTRAN workflow results"
    )
    parser.add_argument(
        "--run-dir", "-r", required=True,
        help="Path to workflow run directory"
    )
    parser.add_argument(
        "--output", "-o", default="summary_slides.pptx",
        help="Output .pptx filename (relative to run-dir)"
    )
    args = parser.parse_args()

    run_dir = Path(args.run_dir)
    if not run_dir.exists():
        print(f"❌ Run directory not found: {run_dir}")
        return

    if not os.getenv("PNNL_API_KEY"):
        print("❌ PNNL_API_KEY environment variable not set")
        return

    data = load_run_data(run_dir)
    slide_content = write_slide_content(data["analysis"])
    output_path = run_dir / args.output
    assemble_slides(slide_content, data["figures"], str(output_path))


if __name__ == "__main__":
    main()